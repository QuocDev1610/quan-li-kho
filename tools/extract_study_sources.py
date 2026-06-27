from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pypdf import PdfReader


def extract_pptx(pptx_path: Path, output_dir: Path) -> None:
    slides_dir = output_dir / "pptx_slides"
    media_dir = output_dir / "pptx_media"
    slides_dir.mkdir(parents=True, exist_ok=True)
    media_dir.mkdir(parents=True, exist_ok=True)

    records: list[dict[str, object]] = []
    presentation = Presentation(pptx_path)
    media_map: list[dict[str, object]] = []
    rel_ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    draw_ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    with zipfile.ZipFile(pptx_path) as archive:
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                shape_text = shape.text.strip()
                if shape_text:
                    texts.append(shape_text)
            record = {"slide": index, "text": texts}
            records.append(record)
            (slides_dir / f"slide-{index:03d}.txt").write_text(
                "\n".join(texts), encoding="utf-8"
            )

            rel_name = f"ppt/slides/_rels/slide{index}.xml.rels"
            slide_name = f"ppt/slides/slide{index}.xml"
            if rel_name in archive.namelist():
                rel_root = ET.fromstring(archive.read(rel_name))
                relationships = {
                    node.attrib["Id"]: node.attrib["Target"]
                    for node in rel_root.findall("r:Relationship", rel_ns)
                }
                slide_root = ET.fromstring(archive.read(slide_name))
                targets = []
                for blip in slide_root.findall(".//a:blip", draw_ns):
                    rel_id = blip.attrib.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if rel_id and rel_id in relationships:
                        targets.append(Path(relationships[rel_id]).name)
                if targets:
                    media_map.append({"slide": index, "media": targets})

        for name in archive.namelist():
            if not name.startswith("ppt/media/") or name.endswith("/"):
                continue
            target = media_dir / Path(name).name
            target.write_bytes(archive.read(name))

    (output_dir / "pptx_text.json").write_text(
        json.dumps(records, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (output_dir / "pptx_media_map.json").write_text(
        json.dumps(media_map, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    with (output_dir / "pptx_text.txt").open("w", encoding="utf-8") as handle:
        for record in records:
            handle.write(f"\n===== SLIDE {record['slide']} =====\n")
            handle.write("\n".join(record["text"]))
            handle.write("\n")


def extract_pdf(pdf_path: Path, output_path: Path) -> None:
    reader = PdfReader(pdf_path)
    with output_path.open("w", encoding="utf-8") as handle:
        for index, page in enumerate(reader.pages, start=1):
            handle.write(f"\n===== PAGE {index} =====\n")
            handle.write(page.extract_text() or "")
            handle.write("\n")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--matrix-pdf", type=Path, required=True)
    parser.add_argument("--answer-pdf", type=Path, required=True)
    parser.add_argument("--output-dir", type=Path, required=True)
    args = parser.parse_args()

    args.output_dir.mkdir(parents=True, exist_ok=True)
    extract_pptx(args.pptx, args.output_dir)
    extract_pdf(args.matrix_pdf, args.output_dir / "matrix_text.txt")
    extract_pdf(args.answer_pdf, args.output_dir / "answer_embedded_text.txt")


if __name__ == "__main__":
    main()
